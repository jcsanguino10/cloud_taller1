import os.path
import logging
# Word - ODT
from docx import Document
import pypandoc

# Excel
from openpyxl import load_workbook
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table

# PPT
from pptx import Presentation
from fpdf import FPDF


logging.basicConfig(level=logging.DEBUG)#Dejo esto por si hay que revisar algo

def conversionWordODTTToPDF(inputfile: str) -> str:
    """
    This function converts a word document (DOCX or ODT) to a PDF file.

    Args:
        inputfile (str): The path to the input file.

    Returns:
        str: The path to the output PDF file.

    Raises:
        Exception: If the input file format is not supported.
    """
    try:
        file_path = os.path.join(os.path.dirname(os.path.abspath(inputfile)),inputfile.split("/")[-1] )
        nameSaveFile = file_path.replace("uploaded", "converted").replace("docx", "pdf") if file_path.endswith("docx") else file_path.replace("uploaded", "converted").replace("odt", "pdf")

        if file_path.endswith('.docx'):
            doc = Document(file_path)
            doc.save(nameSaveFile)
        elif file_path.endswith('.odt'):
            pypandoc.convert_file(file_path, 'pdf', outputfile=nameSaveFile)
        else:
            return {"error": "Formato de archivo no compatible"}

        return nameSaveFile
    except Exception as e:
        print(e)
        raise Exception ("Error converting file")
def conversionExcelToPDF(inputfile: str) -> str:
    """
    This function converts an excel file to a PDF file.

    Args:
        inputfile (str): The path to the input file.

    Returns:
        str: The path to the output PDF file.

    Raises:
        Exception: If the input file format is not supported.
    """
    try:
        file_path = os.path.join(os.path.dirname(os.path.abspath(inputfile)),inputfile.split("/")[-1] )
        nameSaveFile = file_path.replace("uploaded", "converted").replace("xlsx", "pdf")
        
        wb = load_workbook(filename=file_path)
        ws = wb.active  

        data = [[cell.value for cell in row] for row in ws.iter_rows()]

        pdf_table = Table(data)

        pdf = SimpleDocTemplate(nameSaveFile, pagesize=letter)
        pdf.build([pdf_table])
        return nameSaveFile
    except Exception as e:
        print(e)
        raise Exception ("Error converting file")

def conversionPPTToPDF(inputfile: str):
    """
    This function converts a PowerPoint presentation (PPTX) to a PDF file.

    Args:
        inputfile (str): The path to the input file.

    Returns:
        str: The path to the output PDF file.

    Raises:
        Exception: If the input file format is not supported.
    """
    try:
        
        file_path = os.path.join(os.path.dirname(os.path.abspath(inputfile)),inputfile.split("/")[-1] )
        nameSaveFile = file_path.replace("uploaded", "converted").replace("pptx", "pdf")
        prs = Presentation(file_path)
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = ""
                    for paragraph in shape.text_frame.paragraphs:
                        text += paragraph.text + "\n"
                    pdf.set_font("Arial", size=12)
                    pdf.multi_cell(0, 10, txt=text)

            if shape.shape_type == 13:  
                image_path = shape.image.anchor.file
                pdf.image(image_path, x=None, y=None, w=200)
        pdf.output(nameSaveFile)

        return nameSaveFile
    except Exception as e:
        print(e)
        raise Exception ("Error converting file")
    
conversionPPTToPDF('./files/ppt.pptx')